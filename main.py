
# main.py — Streamlit front-end for Router -> (Mongo | Policy) handlers
import streamlit as st
from datetime import datetime, timezone
import traceback
import os
import re

from src.Router_gpt import classify_query, RouteType

# Lazy import Mutlimedia to prevent heavy init on Streamlit load
import importlib
Mutlimedia = importlib.import_module("src.Mutlimedia")
load_in_memory_vectorstore = Mutlimedia.load_in_memory_vectorstore
get_cached_vectorstore = Mutlimedia.get_cached_vectorstore
clear_cache = Mutlimedia.clear_cache
cache_status = Mutlimedia.cache_status
policy_handler = Mutlimedia.policy_handler

from langchain_openai import OpenAIEmbeddings
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from langchain.text_splitter import CharacterTextSplitter

# ===========================================
# PAGE CONFIG
# ===========================================
st.set_page_config(page_title="Mongo_RAG Streamlit", page_icon="📚", layout="wide")
st.title("📚 Mongo_RAG - Streamlit Frontend")

# Tabs for Upload and Query
tab1, tab2 = st.tabs(["📂 Upload & Embed", "💬 Query & Route"])

# ===========================================
# HELPER — Process and Embed Files Safely
# ===========================================
def process_and_embed_files(uploaded_files):
    """Extract text, split safely, and embed once into in-memory vectorstore."""
    try:
        if get_cached_vectorstore():
            st.info("Existing in-memory Chroma DB found — clearing old cache.")
            clear_cache()

        all_texts = []
        splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=100)

        for uploaded_file in uploaded_files:
            file_ext = os.path.splitext(uploaded_file.name)[1].lower()
            text = ""

            if file_ext == ".pdf":
                reader = PdfReader(uploaded_file)
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        text += t + "\n"

            elif file_ext == ".pptx":
                prs = Presentation(uploaded_file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"

            elif file_ext == ".docx":
                doc = Document(uploaded_file)
                for para in doc.paragraphs:
                    text += para.text + "\n"

            # Clean text and avoid large single-blocks
            text = re.sub(r"\s+", " ", text).strip()
            if not text:
                continue

            chunks = splitter.split_text(text)

            # Enforce max chunk length for token safety
            safe_chunks = [c[:1000] if len(c) > 1000 else c for c in chunks]
            all_texts.extend(safe_chunks)

        if not all_texts:
            st.error("❌ No valid text found in uploaded files.")
            return False

        embeddings = OpenAIEmbeddings()
        load_in_memory_vectorstore(embeddings, all_texts, collection_name="ram_store")
        st.success(f"✅ Successfully embedded {len(all_texts)} text chunks in memory!")
        return True

    except Exception as e:
        st.error(f"❌ Embedding process failed: {e}")
        st.code(traceback.format_exc())
        return False


# ===========================================
# TAB 1 — UPLOAD & EMBED
# ===========================================
with tab1:
    st.markdown("""
    ### 📁 Upload Policy Documents
    Upload **PDF**, **DOCX**, or **PPTX** files here to embed them in-memory.
    These embeddings will later be used for answering Policy-related queries.
    """)

    uploaded_files = st.file_uploader(
        "Upload one or more files",
        type=["pdf", "docx", "pptx"],
        accept_multiple_files=True
    )

    if uploaded_files:
        with st.spinner("🔄 Processing and embedding uploaded files..."):
            success = process_and_embed_files(uploaded_files)
        if success:
            st.caption(f"🧮 Cache status: {cache_status()}")
    if st.button("🧹 Clear Cache"):
        msg = clear_cache()
        st.info(msg)

# ===========================================
# TAB 2 — QUERY + ROUTING
# ===========================================
with tab2:
    st.markdown("""
    ### 🧠 Intelligent Query Router
    Enter a natural-language query.  
    The router will:
    1. Classify it as **Policy**, **Document**, or **Both**
    2. Route to appropriate handler(s)
    3. Return detailed results + full trace logs
    """)

    query = st.text_input("Enter your question", placeholder="e.g. What is the leave encashment policy?")

    col1, col2 = st.columns([3, 1])
    with col2:
        if st.button("Run Query"):
            if not query or not query.strip():
                st.warning("⚠️ Please enter a question.")
            else:
                st.session_state['run_time'] = datetime.now(timezone.utc).isoformat()
                st.session_state['query'] = query.strip()

    if 'run_time' in st.session_state:
        st.caption(f"Last run UTC: {st.session_state['run_time']}")

    # ==============================
    # EXECUTION PIPELINE
    # ==============================
    if 'query' in st.session_state and st.session_state['query']:
        q = st.session_state['query']

        st.divider()
        st.subheader("🔍 Step 1 — Classifying Query")

        with st.spinner("Classifying query via Router..."):
            try:
                router_response = classify_query(q)
            except Exception as e:
                st.error(f"❌ Router crashed: {e}")
                st.code(traceback.format_exc())
                st.stop()

        # Handle router output
        if isinstance(router_response, dict) and "error" in router_response:
            st.error(f"❌ Router error at stage **{router_response.get('stage')}**")
            st.code(router_response["error"])
            st.stop()
        else:
            try:
                route, confidence, reason, doc_q, pol_q = router_response
            except Exception as e:
                st.error(f"❌ Unexpected Router output format: {e}")
                st.write(router_response)
                st.stop()

        with st.expander("Router Output (Raw)", expanded=False):
            st.json({
                "route": getattr(route, "value", str(route)),
                "confidence": confidence,
                "reason": reason,
                "doc_query": doc_q,
                "policy_query": pol_q,
            })

        st.success(f"✅ Router classified as **{getattr(route, 'value', str(route)).upper()}** (confidence {confidence})")

        result_text = ""
        stage_logs = []

        # ==============================
        # DOCUMENT HANDLER
        # ==============================
        if getattr(route, "value", str(route)).lower() == "document":
            st.info("🗂️ Running Mongo (Document) handler...")
            stage_logs.append("Router → Document Handler (Mongo)")
            try:
                from src.Mongo import query_mongo
                with st.spinner("Executing Mongo pipeline..."):
                    result_text = query_mongo(q)
                if isinstance(result_text, str) and result_text.startswith("ERROR"):
                    st.error(result_text)
                else:
                    st.success("✅ Document handler executed successfully.")
            except Exception as e:
                stage_logs.append(f"Mongo Handler Error: {e}")
                st.error(f"❌ Document handler failed: {e}")
                st.code(traceback.format_exc())

        # ==============================
        # POLICY HANDLER
        # ==============================
        elif getattr(route, "value", str(route)).lower() == "policy":
            st.info("📜 Running Policy handler...")
            stage_logs.append("Router → Policy Handler")
            try:
                if not get_cached_vectorstore():
                    st.warning("⚠️ No in-memory Chroma DB found — please upload and embed documents first.")
                    st.stop()

                with st.spinner("Executing Policy RAG pipeline..."):
                    result_text = policy_handler(q)

                if not result_text:
                    st.warning("⚠️ No response generated by Policy handler.")
                elif isinstance(result_text, str) and result_text.startswith("ERROR"):
                    st.error(result_text)
                else:
                    st.success("✅ Policy handler executed successfully.")

            except Exception as e:
                stage_logs.append(f"Policy Handler Error: {e}")
                st.error("❌ Policy handler failed:")
                st.code(traceback.format_exc())

        # ==============================
        # BOTH HANDLERS
        # ==============================
        elif getattr(route, "value", str(route)).lower() == "both":
            st.info("🔄 Running both Mongo & Policy handlers...")
            stage_logs.append("Router → Both (Document + Policy)")
            try:
                from src.Mongo import query_mongo
                if not get_cached_vectorstore():
                    st.warning("⚠️ No in-memory Chroma DB found — please upload and embed documents first.")
                    st.stop()

                with st.spinner("Executing Document handler..."):
                    res_doc = query_mongo(doc_q or q)
                with st.spinner("Executing Policy handler..."):
                    res_pol = policy_handler(pol_q or q)

                if any(isinstance(x, str) and x.startswith("ERROR") for x in [res_doc, res_pol]):
                    st.error("⚠️ One or both handlers reported an error.")
                else:
                    st.success("✅ Both handlers executed successfully.")

                result_text = f"--- DOCUMENT RESULT ---\n{res_doc}\n\n--- POLICY RESULT ---\n{res_pol}"

            except Exception as e:
                stage_logs.append(f"Both Handler Error: {e}")
                st.error("❌ Combined handler failure:")
                st.code(traceback.format_exc())

        # ==============================
        # UNKNOWN ROUTE
        # ==============================
        else:
            st.warning(f"⚠️ Unknown route type: {route}")
            stage_logs.append("Unknown Route")

        # ==============================
        # OUTPUT DISPLAY
        # ==============================
        st.divider()
        st.subheader("🧾 Final Result")
        st.code(result_text or "No output generated.", language="text")

        st.subheader("📋 Execution Log")
        for log_entry in stage_logs:
            st.write("-", log_entry)

        st.caption("✅ Process completed successfully — all stages ran in same thread.")
